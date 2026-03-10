from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# Durham brand colors
DURHAM_BLUE = RGBColor(0x00, 0x3F, 0x87)   # deep Durham blue
DURHAM_GREEN = RGBColor(0x00, 0x7A, 0x3D)  # Durham green
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF2, 0xF2, 0xF2)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
MID_GRAY = RGBColor(0x66, 0x66, 0x66)
ACCENT = RGBColor(0x00, 0x9B, 0xDE)  # lighter blue accent

LOGO = "/home/robby-the-bobot/.openclaw/workspace/projects/bms-project-pulse/durham-logo.png"

prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)

def blank_slide(prs):
    blank_layout = prs.slide_layouts[6]
    return prs.slides.add_slide(blank_layout)

def fill_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_textbox(slide, text, left, top, width, height,
                font_size=24, bold=False, color=WHITE,
                align=PP_ALIGN.LEFT, italic=False):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox

def add_rect(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_logo(slide, left, top, height):
    slide.shapes.add_picture(LOGO, left, top, height=height)

def add_bullet_box(slide, title, bullets, left, top, width, height,
                   title_color=DURHAM_BLUE, bullet_color=DARK_GRAY,
                   bg_color=None):
    if bg_color:
        add_rect(slide, left, top, width, height, bg_color)
    # Title
    add_textbox(slide, title, left + Inches(0.15), top + Inches(0.1),
                width - Inches(0.3), Inches(0.5),
                font_size=18, bold=True, color=title_color)
    # Bullets
    tb = slide.shapes.add_textbox(left + Inches(0.15), top + Inches(0.55),
                                   width - Inches(0.3), height - Inches(0.65))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, b in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(4)
        run = p.add_run()
        run.text = f"• {b}"
        run.font.size = Pt(14)
        run.font.color.rgb = bullet_color

# ─────────────────────────────────────────────
# SLIDE 1 — Title Slide
# ─────────────────────────────────────────────
s1 = blank_slide(prs)
fill_bg(s1, DURHAM_BLUE)

# Bottom accent bar
add_rect(s1, 0, Inches(6.8), Inches(13.33), Inches(0.7), DURHAM_GREEN)

# Logo top-left
add_logo(s1, Inches(0.4), Inches(0.3), Inches(1.4))

# Main title
add_textbox(s1, "Let's Talk AI:", Inches(0.5), Inches(2.0),
            Inches(12), Inches(1.0), font_size=52, bold=True,
            color=WHITE, align=PP_ALIGN.CENTER)
add_textbox(s1, "Working Smarter in BMS",
            Inches(0.5), Inches(2.9), Inches(12), Inches(0.9),
            font_size=38, bold=False, color=ACCENT, align=PP_ALIGN.CENTER)

# Subtitle
add_textbox(s1, "Budget & Management Services  |  City of Durham",
            Inches(0.5), Inches(4.1), Inches(12), Inches(0.5),
            font_size=18, color=WHITE, align=PP_ALIGN.CENTER, italic=True)

add_textbox(s1, "March 2026",
            Inches(0.5), Inches(4.7), Inches(12), Inches(0.4),
            font_size=14, color=ACCENT, align=PP_ALIGN.CENTER)

# ─────────────────────────────────────────────
# SLIDE 2 — Icebreaker
# ─────────────────────────────────────────────
s2 = blank_slide(prs)
fill_bg(s2, LIGHT_GRAY)
add_rect(s2, 0, 0, Inches(13.33), Inches(1.0), DURHAM_BLUE)
add_logo(s2, Inches(0.2), Inches(0.12), Inches(0.75))
add_textbox(s2, "Let's Start Here", Inches(1.2), Inches(0.1),
            Inches(10), Inches(0.8), font_size=28, bold=True, color=WHITE)

# Big question
add_rect(s2, Inches(1.0), Inches(1.4), Inches(11.33), Inches(2.8), DURHAM_BLUE)
add_textbox(s2,
    '"Think of a small daily annoyance in your life.\nIs there an app that fixes it?\nAnd if not — what would that app do?"',
    Inches(1.3), Inches(1.6), Inches(10.7), Inches(2.4),
    font_size=26, bold=False, color=WHITE, align=PP_ALIGN.CENTER, italic=True)

add_textbox(s2, "Go around the room — everyone answers.",
            Inches(1.0), Inches(4.5), Inches(11), Inches(0.6),
            font_size=18, color=MID_GRAY, align=PP_ALIGN.CENTER, italic=True)

# ─────────────────────────────────────────────
# SLIDE 3 — What is NotebookLM?
# ─────────────────────────────────────────────
s3 = blank_slide(prs)
fill_bg(s3, WHITE)
add_rect(s3, 0, 0, Inches(13.33), Inches(1.0), DURHAM_BLUE)
add_logo(s3, Inches(0.2), Inches(0.12), Inches(0.75))
add_textbox(s3, "What is NotebookLM?", Inches(1.2), Inches(0.1),
            Inches(10), Inches(0.8), font_size=28, bold=True, color=WHITE)

add_textbox(s3,
    "Google's AI research assistant — but it only knows what YOU give it.\nNo hallucinating random internet stuff. Just your documents, made searchable and smart.",
    Inches(0.5), Inches(1.15), Inches(12.3), Inches(0.9),
    font_size=16, color=DARK_GRAY, italic=True)

# Two columns
add_bullet_box(s3, "📂  What You Can Feed It",
    ["PDFs, Google Docs, Slides", "YouTube video links", "Website URLs", "Copy-pasted text", "Up to 50 sources per notebook (free)"],
    Inches(0.4), Inches(2.2), Inches(5.8), Inches(3.8), bg_color=LIGHT_GRAY,
    title_color=DURHAM_BLUE, bullet_color=DARK_GRAY)

add_bullet_box(s3, "⚡  What It Does",
    ["Chat with your docs — get answers with citations", "Auto-summary of everything in the notebook", "Audio Overview — AI podcast from your content 🎙️", "Slide Deck generator (new!)", "Study guides, FAQs, briefing docs"],
    Inches(6.6), Inches(2.2), Inches(6.3), Inches(3.8), bg_color=LIGHT_GRAY,
    title_color=DURHAM_BLUE, bullet_color=DARK_GRAY)

add_rect(s3, 0, Inches(7.2), Inches(13.33), Inches(0.3), DURHAM_GREEN)

# ─────────────────────────────────────────────
# SLIDE 4 — NotebookLM Best Uses
# ─────────────────────────────────────────────
s4 = blank_slide(prs)
fill_bg(s4, WHITE)
add_rect(s4, 0, 0, Inches(13.33), Inches(1.0), DURHAM_BLUE)
add_logo(s4, Inches(0.2), Inches(0.12), Inches(0.75))
add_textbox(s4, "NotebookLM: Best Use Cases", Inches(1.2), Inches(0.1),
            Inches(10), Inches(0.8), font_size=28, bold=True, color=WHITE)

use_cases = [
    ("📋", "Digesting Long Reports", "Upload a 200-page budget doc and ask it plain-English questions in seconds."),
    ("🎥", "Meeting Recordings & Videos", "Drop in a YouTube link — it reads transcripts and answers questions about the content."),
    ("🔍", "Research Without Reading Everything", "Throw in 10 sources and ask it to synthesize across all of them."),
    ("🗂️", "Building a Knowledge Base", "Centralize scattered docs into one smart, searchable notebook."),
]

for i, (emoji, title, desc) in enumerate(use_cases):
    col = i % 2
    row = i // 2
    left = Inches(0.4 + col * 6.5)
    top = Inches(1.2 + row * 2.7)
    add_rect(s4, left, top, Inches(6.1), Inches(2.4), LIGHT_GRAY)
    add_textbox(s4, f"{emoji}  {title}", left + Inches(0.15), top + Inches(0.1),
                Inches(5.8), Inches(0.5), font_size=16, bold=True, color=DURHAM_BLUE)
    add_textbox(s4, desc, left + Inches(0.15), top + Inches(0.6),
                Inches(5.8), Inches(1.6), font_size=13, color=DARK_GRAY)

add_rect(s4, 0, Inches(7.2), Inches(13.33), Inches(0.3), DURHAM_GREEN)

# ─────────────────────────────────────────────
# SLIDE 5 — What is Vibe Coding?
# ─────────────────────────────────────────────
s5 = blank_slide(prs)
fill_bg(s5, DURHAM_BLUE)
add_rect(s5, 0, Inches(6.8), Inches(13.33), Inches(0.7), DURHAM_GREEN)
add_logo(s5, Inches(0.2), Inches(0.12), Inches(0.75))

add_textbox(s5, "What is Vibe Coding? 💻", Inches(0.5), Inches(0.8),
            Inches(12.3), Inches(0.9), font_size=36, bold=True,
            color=WHITE, align=PP_ALIGN.CENTER)

add_textbox(s5,
    "Building software by describing what you want\nin plain English — no technical knowledge required.",
    Inches(0.5), Inches(1.8), Inches(12.3), Inches(0.9),
    font_size=20, color=ACCENT, align=PP_ALIGN.CENTER, italic=True)

# Three boxes
boxes = [
    ("1️⃣", "Describe It", "Tell the AI what you want in plain English"),
    ("2️⃣", "AI Builds It", "Code is written automatically in seconds"),
    ("3️⃣", "Ship It", "Deploy a live app with one command"),
]
for i, (num, title, desc) in enumerate(boxes):
    left = Inches(0.5 + i * 4.2)
    add_rect(s5, left, Inches(3.0), Inches(3.8), Inches(2.5), RGBColor(0x00, 0x2D, 0x6A))
    add_textbox(s5, num, left + Inches(0.1), Inches(3.1),
                Inches(3.6), Inches(0.5), font_size=24, align=PP_ALIGN.CENTER, color=ACCENT)
    add_textbox(s5, title, left + Inches(0.1), Inches(3.65),
                Inches(3.6), Inches(0.5), font_size=18, bold=True,
                color=WHITE, align=PP_ALIGN.CENTER)
    add_textbox(s5, desc, left + Inches(0.1), Inches(4.2),
                Inches(3.6), Inches(1.0), font_size=13, color=ACCENT,
                align=PP_ALIGN.CENTER)

add_textbox(s5, "It's not the future. It's Tuesday.",
            Inches(0.5), Inches(5.8), Inches(12.3), Inches(0.5),
            font_size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER, italic=True)

# ─────────────────────────────────────────────
# SLIDE 6 — The Live Demo Prompt
# ─────────────────────────────────────────────
s6 = blank_slide(prs)
fill_bg(s6, WHITE)
add_rect(s6, 0, 0, Inches(13.33), Inches(1.0), DURHAM_BLUE)
add_logo(s6, Inches(0.2), Inches(0.12), Inches(0.75))
add_textbox(s6, "Let's Build Something. Live. 🚀", Inches(1.2), Inches(0.1),
            Inches(11), Inches(0.8), font_size=28, bold=True, color=WHITE)

add_textbox(s6, "The prompt I'm sending to my AI right now:",
            Inches(0.5), Inches(1.1), Inches(12), Inches(0.5),
            font_size=16, color=MID_GRAY, italic=True)

# Prompt box
add_rect(s6, Inches(0.4), Inches(1.7), Inches(12.5), Inches(3.2), LIGHT_GRAY)
add_textbox(s6,
    '"Hey Robby — build me a simple project tracker app called BMS Project Pulse. '
    'It should have a clean dashboard showing all my projects with a status badge '
    '(On Track, At Risk, or Done), a last-updated note, and the ability to add new projects. '
    'Pre-load it with 5 fake Durham BMS projects so it feels real. '
    'Make it look professional — dark header, clean cards. '
    'Deploy it to GitHub Pages when it\'s done."',
    Inches(0.6), Inches(1.85), Inches(12.1), Inches(2.9),
    font_size=15, color=DARK_GRAY, italic=True)

add_textbox(s6, "⬆️  That's it. Plain English. No code. Let's see what happens...",
            Inches(0.5), Inches(5.1), Inches(12.3), Inches(0.6),
            font_size=16, bold=True, color=DURHAM_BLUE, align=PP_ALIGN.CENTER)

add_textbox(s6, "While it builds → NotebookLM demo",
            Inches(0.5), Inches(5.8), Inches(12.3), Inches(0.5),
            font_size=14, color=MID_GRAY, align=PP_ALIGN.CENTER, italic=True)

add_rect(s6, 0, Inches(7.2), Inches(13.33), Inches(0.3), DURHAM_GREEN)

# ─────────────────────────────────────────────
# SLIDE 7 — AI Isn't Replacing You
# ─────────────────────────────────────────────
s7 = blank_slide(prs)
fill_bg(s7, DURHAM_BLUE)
add_rect(s7, 0, Inches(6.8), Inches(13.33), Inches(0.7), DURHAM_GREEN)
add_logo(s7, Inches(0.2), Inches(0.12), Inches(0.75))

add_textbox(s7, "AI Isn't Coming for Your Job.", Inches(0.5), Inches(0.9),
            Inches(12.3), Inches(0.9), font_size=36, bold=True,
            color=WHITE, align=PP_ALIGN.CENTER)

add_textbox(s7,
    "The people who use AI well will outperform the people who don't.\nYou're not being replaced. You're getting a superpower.",
    Inches(0.5), Inches(1.85), Inches(12.3), Inches(0.9),
    font_size=18, color=ACCENT, align=PP_ALIGN.CENTER, italic=True)

# Two columns
add_rect(s7, Inches(0.4), Inches(2.9), Inches(5.9), Inches(3.5), RGBColor(0x00, 0x2D, 0x6A))
add_textbox(s7, "✅  AI Is Great At:",
            Inches(0.6), Inches(3.0), Inches(5.5), Inches(0.5),
            font_size=16, bold=True, color=ACCENT)
goods = ["Digesting long documents instantly", "Repetitive drafting & formatting", "Building tools you'd wait months for", "Finding patterns in data"]
for i, g in enumerate(goods):
    add_textbox(s7, f"• {g}", Inches(0.6), Inches(3.6 + i*0.55), Inches(5.5), Inches(0.5),
                font_size=13, color=WHITE)

add_rect(s7, Inches(7.0), Inches(2.9), Inches(5.9), Inches(3.5), RGBColor(0x00, 0x2D, 0x6A))
add_textbox(s7, "❌  AI Can't Replace:",
            Inches(7.2), Inches(3.0), Inches(5.5), Inches(0.5),
            font_size=16, bold=True, color=ACCENT)
bads = ["Your department knowledge", "Relationships & trust", "Political judgment", "Accountability — that's still you"]
for i, b in enumerate(bads):
    add_textbox(s7, f"• {b}", Inches(7.2), Inches(3.6 + i*0.55), Inches(5.5), Inches(0.5),
                font_size=13, color=WHITE)

# ─────────────────────────────────────────────
# SLIDE 8 — What AI Is Good At
# ─────────────────────────────────────────────
s8 = blank_slide(prs)
fill_bg(s8, WHITE)
add_rect(s8, 0, 0, Inches(13.33), Inches(1.0), DURHAM_BLUE)
add_logo(s8, Inches(0.2), Inches(0.12), Inches(0.75))
add_textbox(s8, "Start Small. Start Now.", Inches(1.2), Inches(0.1),
            Inches(11), Inches(0.8), font_size=28, bold=True, color=WHITE)

add_textbox(s8,
    "You don't need to be technical. You don't need a budget.\nYou need a question and a keyboard.",
    Inches(0.5), Inches(1.15), Inches(12.3), Inches(0.8),
    font_size=18, color=DARK_GRAY, italic=True, align=PP_ALIGN.CENTER)

questions = [
    ("📝", "What's one report you dread writing every week?"),
    ("📅", "What's one meeting you spend 30 minutes prepping for?"),
    ("🛠️", "What's one tool you wish existed?"),
]
for i, (emoji, q) in enumerate(questions):
    left = Inches(0.4 + i * 4.2)
    add_rect(s8, left, Inches(2.2), Inches(3.9), Inches(2.5), LIGHT_GRAY)
    add_textbox(s8, emoji, left + Inches(0.1), Inches(2.3),
                Inches(3.7), Inches(0.6), font_size=30, align=PP_ALIGN.CENTER, color=DURHAM_BLUE)
    add_textbox(s8, q, left + Inches(0.1), Inches(2.95),
                Inches(3.7), Inches(1.5), font_size=14, color=DARK_GRAY,
                align=PP_ALIGN.CENTER)

add_rect(s8, Inches(0.4), Inches(5.0), Inches(12.5), Inches(1.4), DURHAM_BLUE)
add_textbox(s8, "That's where AI starts.",
            Inches(0.5), Inches(5.15), Inches(12.1), Inches(0.5),
            font_size=24, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_textbox(s8, "And today, you've already seen it work.",
            Inches(0.5), Inches(5.65), Inches(12.1), Inches(0.5),
            font_size=16, color=ACCENT, align=PP_ALIGN.CENTER, italic=True)

add_rect(s8, 0, Inches(7.2), Inches(13.33), Inches(0.3), DURHAM_GREEN)

# ─────────────────────────────────────────────
# Save
# ─────────────────────────────────────────────
out = "/home/robby-the-bobot/.openclaw/workspace/projects/bms-project-pulse/BMS_AI_Presentation.pptx"
prs.save(out)
print(f"Saved: {out}")
